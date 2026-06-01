"""Edit .pptx on disk with python-pptx (close file in PowerPoint first if locked)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation


def add_slide(pptx_path: Path, layout_name: str) -> str:
    """Append a new slide using a named layout from the deck's slide master."""
    prs = Presentation(str(pptx_path))
    layout = None
    for sl in prs.slide_layouts:
        if sl.name.lower() == layout_name.lower():
            layout = sl
            break
    if layout is None:
        available = [sl.name for sl in prs.slide_layouts]
        raise ValueError(f"Layout {layout_name!r} not found. Available: {available}")
    prs.slides.add_slide(layout)
    prs.save(str(pptx_path))
    return f"Added slide with layout '{layout.name}' (now {len(prs.slides)} slides total)."


def populate_plain_text(
    pptx_path: Path,
    slide_number: int,
    placeholder_name: str,
    content: str,
) -> str:
    """
    Set plain text on the first shape whose name contains placeholder_name (case-insensitive).

    v0 does not parse HTML/LaTeX; strips tags crudely for display text only.
    """
    prs = Presentation(str(pptx_path))
    if slide_number < 1 or slide_number > len(prs.slides):
        raise ValueError(f"slide_number out of range: {slide_number} (deck has {len(prs.slides)} slides)")
    slide = prs.slides[slide_number - 1]
    needle = placeholder_name.lower()
    updated = 0
    for shape in slide.shapes:
        name = (getattr(shape, "name", "") or "").lower()
        if needle in name and getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
            shape.text_frame.text = _strip_simple_tags(content)
            updated += 1
    if updated == 0:
        for ph in slide.placeholders:
            name = (getattr(ph, "name", "") or "").lower()
            if needle in name and ph.has_text_frame:
                ph.text_frame.text = _strip_simple_tags(content)
                updated += 1
    if updated == 0:
        raise ValueError(
            f"No placeholder/shape matched {placeholder_name!r} on slide {slide_number}. "
            "Use slide_snapshot (when available) or inspect shape names in PowerPoint."
        )
    prs.save(str(pptx_path))
    return f"Updated {updated} shape(s) on slide {slide_number}."


def _strip_simple_tags(s: str) -> str:
    """Very small v0 stripper; not HTML-safe."""
    out = s
    for tag in ("<b>", "</b>", "<i>", "</i>", "<u>", "</u>"):
        out = out.replace(tag, "")
    return out
